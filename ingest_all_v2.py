import chromadb
import ollama
import uuid
import os
import io
import sys
import datetime
import time
import json
import fitz

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

from config_loader import get_config
from hybrid_retrieval import build_hybrid_metadata

from docx import Document
from pptx import Presentation
from config_loader import get_config

# ==================== 加载配置 ====================
config = get_config()

DB_PATH = config.db_path
PDF_FOLDER = config.pdf_folder
WORD_FOLDER = config.word_folder
PPT_FOLDER = config.ppt_folder
IMG_FOLDER = config.img_folder
COLLECTION_NAME = config.collection_name
CHUNK_SIZE = config.chunk_size
OVERLAP = config.overlap
BATCH_SIZE = config.batch_size
EMBEDDING_MODEL = config.embedding_model
OCR_LANGUAGES = config.ocr_languages
CHROMA_MODE = config.chroma_mode
CHROMA_HOST = config.chroma_server_host
CHROMA_PORT = config.chroma_server_port

# 支持的图片格式
IMG_EXTS = ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.gif', '.webp')
# =================================================

# 根据配置选择连接方式
if CHROMA_MODE == 'remote':
    client = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
    print(f"🔗 远程模式: 连接 {CHROMA_HOST}:{CHROMA_PORT}")
else:
    client = chromadb.PersistentClient(path=DB_PATH)
    print(f"💾 本地模式: {DB_PATH}")

col = client.get_or_create_collection(
    name=COLLECTION_NAME,
    metadata={"hnsw:space": "cosine"}
)

def check_dependencies():
    """检查所有依赖是否就绪"""
    errors = []
    warnings = []
    
    # 检查 Ollama 服务
    try:
        ollama.list()
    except Exception as e:
        errors.append(f"Ollama 服务未运行: {e}")
    
    # 检查 Tesseract OCR（图片非必要时仅警告）
    has_img = os.path.exists(IMG_FOLDER) and any(
        f.lower().endswith(IMG_EXTS) for f in os.listdir(IMG_FOLDER)
    ) if os.path.exists(IMG_FOLDER) else False
    try:
        pytesseract.get_tesseract_version()
    except Exception as e:
        if has_img:
            errors.append(f"Tesseract OCR 未安装或配置错误: {e}")
        else:
            print(f"  [注意] Tesseract 未安装（无图片文件，可忽略）")
    
    # 检查文件夹(至少需要一个存在)
    folders_exist = False
    for folder, name in [(PDF_FOLDER, 'PDF'), (WORD_FOLDER, 'Word'), 
                         (PPT_FOLDER, 'PPT'), (IMG_FOLDER, '图片')]:
        if os.path.exists(folder):
            folders_exist = True
        else:
            warnings.append(f"{name} 文件夹不存在: {folder}")
    
    if not folders_exist:
        errors.append("所有文档文件夹都不存在")
    
    if errors:
        print("❌ 依赖检查失败:")
        for err in errors:
            print(f"  - {err}")
        return False
    
    if warnings:
        print("⚠️  警告:")
        for warn in warnings:
            print(f"  - {warn}")
    
    print("✅ 依赖检查通过")
    return True

def embed_batch(texts: list) -> list:
    """批量生成嵌入向量,提高性能"""
    embeddings = []
    for idx, text in enumerate(texts):
        try:
            # nomic-embed-text-v1 最大输入长度为 512 tokens
            # 保守估计: 1个中文字符 ≈ 1.5 tokens, 1个英文单词 ≈ 1.3 tokens
            # 设置安全限制为 450 字符
            safe_text = text[:450]
            
            if len(text) > 450:
                print(f"  [警告] 文本过长 ({len(text)} 字符),已截断至 450 字符")
            
            r = ollama.embeddings(model=EMBEDDING_MODEL, prompt=safe_text)
            embeddings.append(r["embedding"])
            # 避免请求过快,每5个请求延迟0.1秒
            if (idx + 1) % 5 == 0:
                time.sleep(0.1)
        except Exception as e:
            print(f"  [错误] 嵌入生成失败: {e}")
            embeddings.append(None)
    return embeddings

def split_by_separators(text: str, separators: list) -> list:
    """按分隔符列表递归分割文本"""
    for separator in separators:
        if separator in text:
            parts = text.split(separator)
            # 过滤空字符串并保留分隔符
            result = []
            for i, part in enumerate(parts):
                if part.strip():
                    result.append(part)
            if len(result) > 1:
                return result
    return [text]

def chunk_text_recursive(text: str, chunk_size: int = None, chunk_overlap: int = None):
    """
    递归式文本分块(Recursive Chunking)
    
    策略:
    1. 先尝试按大边界分割(段落、双换行)
    2. 如果片段仍太大,按小边界分割(句子、单换行)
    3. 最后才按字符数硬切分
    
    Args:
        text: 待分块文本
        chunk_size: 最大块大小(默认使用配置)
        chunk_overlap: 重叠大小(默认使用配置)
    
    Returns:
        分块后的文本列表
    """
    if chunk_size is None:
        chunk_size = CHUNK_SIZE
    if chunk_overlap is None:
        chunk_overlap = OVERLAP
    
    if not text or len(text) < 50:
        return []
    
    # 如果文本长度小于chunk_size,直接返回
    if len(text) <= chunk_size:
        return [text]
    
    # 定义分割优先级(从大到小)
    separators = [
        '\n\n',      # 段落分隔
        '\n',        # 换行
        '。',        # 中文句号
        '！',        # 中文感叹号
        '？',        # 中文问号
        '. ',        # 英文句号+空格
        '! ',        # 英文感叹号+空格
        '? ',        # 英文问号+空格
        '；',        # 中文分号
        '; ',        # 英文分号+空格
        '，',        # 中文逗号
        ', ',        # 英文逗号+空格
    ]
    
    # 递归分割
    chunks = _recursive_split(text, chunk_size, chunk_overlap, separators, 0)
    
    # 合并过小的块(可选优化)
    chunks = merge_small_chunks(chunks, chunk_size * 0.8)
    
    return chunks

def _recursive_split(text: str, chunk_size: int, chunk_overlap: int, 
                     separators: list, level: int) -> list:
    """递归分割内部函数"""
    # 如果文本已经足够小,直接返回
    if len(text) <= chunk_size:
        return [text]
    
    # 如果已经尝试完所有分隔符,按字符硬切分
    if level >= len(separators):
        return hard_chunk(text, chunk_size, chunk_overlap)
    
    separator = separators[level]
    
    # 尝试用当前分隔符分割
    if separator in text:
        parts = text.split(separator)
        
        # 重新组合,保持语义完整
        chunks = []
        current_chunk = ""
        
        for part in parts:
            part = part.strip()
            if not part:
                continue
            
            # 添加分隔符(除了最后一个)
            if current_chunk:
                test_chunk = current_chunk + separator + part
            else:
                test_chunk = part
            
            # 如果加入后超过限制,保存当前块,开始新块
            if len(test_chunk) > chunk_size and current_chunk:
                chunks.append(current_chunk)
                current_chunk = part
            else:
                current_chunk = test_chunk
        
        # 添加最后一个块
        if current_chunk:
            chunks.append(current_chunk)
        
        # 如果分割后仍有大块,继续递归分割
        final_chunks = []
        for chunk in chunks:
            if len(chunk) > chunk_size:
                # 递归到下一级分隔符
                sub_chunks = _recursive_split(chunk, chunk_size, chunk_overlap, 
                                             separators, level + 1)
                final_chunks.extend(sub_chunks)
            else:
                final_chunks.append(chunk)
        
        return final_chunks
    else:
        # 当前分隔符不存在,尝试下一级
        return _recursive_split(text, chunk_size, chunk_overlap, 
                               separators, level + 1)

def hard_chunk(text: str, chunk_size: int, chunk_overlap: int) -> list:
    """硬切分(最后手段)"""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - chunk_overlap
    return chunks

def merge_small_chunks(chunks: list, min_size: float) -> list:
    """合并过小的相邻块"""
    if not chunks:
        return chunks
    
    merged = [chunks[0]]
    
    for i in range(1, len(chunks)):
        if len(merged[-1]) < min_size:
            # 合并到前一个块
            merged[-1] = merged[-1] + "\n\n" + chunks[i]
        else:
            merged.append(chunks[i])
    
    return merged

def chunk_text(text: str):
    """将文本分块(使用递归式切片)"""
    return chunk_text_recursive(text)

def ocr_image(image_bytes: bytes) -> str:
    """OCR 识别图片文字"""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        text = pytesseract.image_to_string(img, lang=OCR_LANGUAGES)
        return text.strip()
    except Exception as e:
        return f"[OCR失败: {e}]"

def ocr_image_file(img_path: str) -> str:
    """直接读取图片文件OCR"""
    try:
        with open(img_path, 'rb') as f:
            return ocr_image(f.read())
    except Exception as e:
        return f"[OCR失败: {e}]"

def extract_pdf(pdf_path: str) -> str:
    """提取 PDF 内容,包括文字和 OCR 图片"""
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"  [错误] 无法打开 PDF: {e}")
        return ""
    
    all_blocks = []
    processed_xrefs = set()
    
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # 提取文字
            try:
                text = page.get_text()
                if text.strip():
                    all_blocks.append(f"[第{page_num+1}页 文字]\n{text.strip()}")
            except Exception as e:
                print(f"  [警告] 第{page_num+1}页文字提取失败: {e}")
            
            # 提取图片并 OCR
            try:
                images = page.get_images(full=True)
                for img_index, img in enumerate(images, start=1):
                    xref = img[0]
                    if xref in processed_xrefs:
                        continue
                    processed_xrefs.add(xref)
                    
                    try:
                        base_image = doc.extract_image(xref)
                        if not base_image:
                            continue
                        ocr_text = ocr_image(base_image["image"])
                        if ocr_text and len(ocr_text) > 10:
                            all_blocks.append(f"[第{page_num+1}页 图{img_index} OCR]\n{ocr_text}")
                    except Exception as e:
                        print(f"  [警告] 第{page_num+1}页图{img_index} OCR 失败: {e}")
                        continue
            except Exception as e:
                print(f"  [警告] 第{page_num+1}页图片处理失败: {e}")
    finally:
        doc.close()
    
    return "\n\n".join(all_blocks)


def iter_pdf_blocks(pdf_path: str):
    """逐页生成 PDF 文本块，避免一次性加载整个大文件"""
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"  [错误] 无法打开 PDF: {e}")
        return

    processed_xrefs = set()
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]

            try:
                text = page.get_text()
                if text.strip():
                    yield f"[第{page_num+1}页 文字]\n{text.strip()}"
            except Exception as e:
                print(f"  [警告] 第{page_num+1}页文字提取失败: {e}")

            try:
                images = page.get_images(full=True)
                for img_index, img in enumerate(images, start=1):
                    xref = img[0]
                    if xref in processed_xrefs:
                        continue
                    processed_xrefs.add(xref)

                    try:
                        base_image = doc.extract_image(xref)
                        if not base_image:
                            continue
                        ocr_text = ocr_image(base_image["image"])
                        if ocr_text and len(ocr_text) > 10:
                            yield f"[第{page_num+1}页 图{img_index} OCR]\n{ocr_text}"
                    except Exception as e:
                        print(f"  [警告] 第{page_num+1}页图{img_index} OCR 失败: {e}")
                        continue
            except Exception as e:
                print(f"  [警告] 第{page_num+1}页图片处理失败: {e}")
    finally:
        doc.close()


def _add_chunks_to_collection(chunks: list, source_name: str, start_index: int, file_type: str, timestamp: str):
    ids = [str(uuid.uuid4()) for _ in chunks]
    embs = embed_batch(chunks)
    valid_data = []
    for emb, doc_text, id_, chunk_index in zip(embs, chunks, ids, range(start_index, start_index + len(chunks))):
        if emb is not None:
            valid_data.append((emb, doc_text, id_, chunk_index))
    if not valid_data:
        return 0
    valid_embs, valid_docs, valid_ids, valid_indexes = zip(*valid_data)
    meta_list = []
    for doc_text, idx in zip(valid_docs, valid_indexes):
        base_meta = {
            "source": source_name,
            "type": file_type,
            "chunk_index": idx,
            "timestamp": timestamp,
        }
        meta_list.append(base_meta)

    col.add(
        embeddings=list(valid_embs),
        documents=list(valid_docs),
        ids=list(valid_ids),
        metadatas=meta_list,
    )
    return len(valid_docs)


def ingest_pdf_file(pdf_path: str, source_name: str, timestamp: str):
    print(f"处理: {source_name} ...")
    existing = col.get(where={"source": source_name})
    if existing["ids"]:
        print(f"  [跳过] 已存在 {len(existing['ids'])} 条记录")
        return 0

    file_type = 'pdf'
    total = 0
    chunk_index = 0
    batch = []

    for block in iter_pdf_blocks(pdf_path):
        if not block or len(block.strip()) < 20:
            continue
        chunks = chunk_text(block)
        if not chunks:
            continue
        for chunk in chunks:
            batch.append(chunk)
            if len(batch) >= BATCH_SIZE:
                total += _add_chunks_to_collection(batch, source_name, chunk_index, file_type, timestamp)
                chunk_index += len(batch)
                batch = []

    if batch:
        total += _add_chunks_to_collection(batch, source_name, chunk_index, file_type, timestamp)
        chunk_index += len(batch)

    if chunk_index == 0:
        print(f"  [跳过] 内容过少")
        return 0

    print(f"  已生成 {chunk_index} 段, 已入库 {total} 条记录")
    return total


def extract_docx(docx_path: str) -> str:
    """提取 Word 文档（含段落+表格）"""
    try:
        doc = Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    tables_text.append(" | ".join(row_text))
        all_text = paragraphs + tables_text
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"  [错误] 无法打开 Word: {e}")
        return ""

def extract_pptx(pptx_path: str) -> str:
    """提取 PPT 文字（逐页）"""
    try:
        prs = Presentation(pptx_path)
        all_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[第{i}页]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if len(slide_texts) > 1:
                all_text.append("\n".join(slide_texts))
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"  [错误] 无法打开 PPT: {e}")
        return ""

def extract_file(file_path: str, source_name: str) -> str:
    """按后缀分发提取器"""
    ext = os.path.splitext(source_name)[1].lower()
    if ext == '.pdf':
        return extract_pdf(file_path)
    elif ext == '.docx':
        return extract_docx(file_path)
    elif ext == '.pptx':
        return extract_pptx(file_path)
    elif ext in IMG_EXTS:
        return ocr_image_file(file_path)
    else:
        print(f"  [跳过] 不支持的格式: {ext}")
        return ""

def ingest_file(file_path: str, source_name: str, timestamp: str):
    """处理单个文件并入库"""
    print(f"处理: {source_name} ...")
    
    # 检查是否已入库
    existing = col.get(where={"source": source_name})
    if existing["ids"]:
        print(f"  [跳过] 已存在 {len(existing['ids'])} 条记录")
        return 0
    
    full_text = extract_file(file_path, source_name)
    if not full_text or len(full_text) < 50:
        print(f"  [跳过] 内容过少或提取失败")
        return 0
    
    # 图片OCR结果通常较短，如果不足chunk_size则整段入库
    ext = os.path.splitext(source_name)[1].lower()
    if ext in IMG_EXTS and len(full_text) < CHUNK_SIZE:
        chunks = [full_text]
    else:
        chunks = chunk_text(full_text)
    
    if not chunks:
        return 0
    
    print(f"  提取 {len(chunks)} 段，开始入库...")
    total = 0
    file_type = ext.replace('.', '') if ext.startswith('.') else 'unknown'
    
    for i in range(0, len(chunks), BATCH_SIZE):
        batch = chunks[i:i+BATCH_SIZE]
        ids = [str(uuid.uuid4()) for _ in batch]
        
        # 使用批量嵌入
        embs = embed_batch(batch)
        
        # 过滤掉失败的嵌入
        valid_data = [
            (emb, doc, id_) 
            for emb, doc, id_ in zip(embs, batch, ids) 
            if emb is not None
        ]
        
        if not valid_data:
            print(f"  [警告] 批次 {i//BATCH_SIZE + 1} 全部嵌入失败")
            continue
        
        valid_embs, valid_docs, valid_ids = zip(*valid_data)

        meta_list = []
        for j, doc_text in enumerate(valid_docs):
            meta = {
                "source": source_name,
                "type": file_type,
                "chunk_index": i + j,
                "timestamp": timestamp,
            }
            meta_list.append(meta)

        col.add(
            embeddings=list(valid_embs),
            documents=list(valid_docs),
            ids=list(valid_ids),
            metadatas=meta_list,
        )
        total += len(valid_docs)
        print(f"  已入库 {total}/{len(chunks)}")
    
    return total

def main():
    """主函数:处理所有支持的文件格式"""
    print("=" * 60)
    print("知识库入库工具 (PDF + Word + PPT + 图片)")
    print("=" * 60)
    
    # 显示当前配置
    print(f"\n{config.get_profile_info()}")
    print("=" * 60)
    
    # 依赖检查
    if not check_dependencies():
        print("\n请先解决上述依赖问题")
        return
    
    # 收集所有文件
    all_files = []
    folders = [
        (PDF_FOLDER, '.pdf', 'PDF'),
        (WORD_FOLDER, '.docx', 'Word'),
        (PPT_FOLDER, '.pptx', 'PPT'),
        (IMG_FOLDER, IMG_EXTS, '图片'),
    ]
    
    for folder, exts, type_name in folders:
        if not os.path.exists(folder):
            print(f"⚠️  目录不存在,跳过: {folder}")
            continue
        
        if isinstance(exts, str):
            exts = (exts,)
        
        files = [(os.path.join(folder, f), f) for f in os.listdir(folder) if f.lower().endswith(exts)]
        all_files.extend(files)
        print(f"📂 {type_name}: {len(files)} 个文件")
    
    if not all_files:
        print("\n❌ 没有可处理的文件")
        return
    
    print(f"\n共发现 {len(all_files)} 个文件，开始处理...")
    print("=" * 60)
    
    # 生成统一的时间戳
    timestamp = datetime.datetime.now().isoformat()
    
    total_chunks = 0
    success_count = 0
    skip_count = 0
    error_count = 0
    
    for idx, (fp, fn) in enumerate(all_files, 1):
        # 显示总体进度
        print(f"\n[{idx}/{len(all_files)}] ", end="")
        
        try:
            n = ingest_file(fp, fn, timestamp)
            total_chunks += n
            if n > 0:
                success_count += 1
                print(f"[完成] {fn} → {n} 段")
            else:
                skip_count += 1
                print(f"[跳过] {fn}")
        except Exception as e:
            error_count += 1
            print(f"[错误] {fn} 处理失败: {e}")
    
    print("\n" + "=" * 60)
    print(f"全部完成!")
    print(f"  成功: {success_count} 个文件")
    print(f"  跳过: {skip_count} 个文件")
    print(f"  错误: {error_count} 个文件")
    print(f"  共入库 {total_chunks} 段")
    print("=" * 60)
    print("\n提示: 等待百度云盘同步完成，再去另一台查询")

if __name__ == "__main__":
    main()
